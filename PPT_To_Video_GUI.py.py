import os
import time
import tkinter as tk
from tkinter import filedialog, messagebox
import customtkinter as ctk
from pptx import Presentation
import win32com.client as win32
from pathlib import Path
import sys
import threading
import asyncio
import edge_tts

# 修复打包时的imageio问题 - 在导入moviepy之前处理
try:
    import imageio
    # 禁用imageio的元数据检查
    os.environ['IMAGEIO_NO_VERSION_CHECK'] = '1'
except Exception:
    pass

# 兼容moviepy的导入 - 放在后面
try:
    import moviepy
except Exception:
    pass

# 兼容moviepy的导入
try:
    from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips
    try:
        from moviepy.editor import TextClip, CompositeVideoClip
    except:
        TextClip = None
        CompositeVideoClip = None
except:
    from moviepy import AudioFileClip, ImageClip, concatenate_videoclips
    TextClip = None
    CompositeVideoClip = None


def get_temp_path():
    if getattr(sys, 'frozen', False):
        base_path = os.path.dirname(sys.executable)
    else:
        base_path = os.path.abspath(".")
    return Path(base_path) / "ppt_temp_files"


ctk.set_appearance_mode("light")
ctk.set_default_color_theme("blue")


class PPTToVideoGUI(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("PPT转带语音视频工具（微软高质量）")
        self.geometry("750x550")
        self.resizable(False, False)

        self.ppt_path = tk.StringVar()
        self.output_name = tk.StringVar(value="输出视频")
        self.subtitle_font_size = tk.IntVar(value=20)
        self.is_converting = False
        self.stop_requested = False

        self.temp = get_temp_path()
        self.temp.mkdir(exist_ok=True)
        (self.temp / "images").mkdir(exist_ok=True)
        (self.temp / "audios").mkdir(exist_ok=True)

        self.create_ui()
        self.log("✅ 软件启动成功！")
        self.log("💡 使用微软Edge高质量在线语音合成")
        self.log("📝 语音内容来自PPT备注，无备注则使用默认文本")
        self.log("📄 自动生成SRT字幕文件，可直接拖入播放器使用")

    def create_ui(self):
        title = ctk.CTkLabel(self, text="PPT 转 带语音视频（微软高质量TTS）", font=("微软雅黑", 20, "bold"))
        title.pack(pady=15)

        frame1 = ctk.CTkFrame(self)
        frame1.pack(padx=25, pady=8, fill="x")
        ctk.CTkLabel(frame1, text="选择PPT文件：", font=("微软雅黑", 12)).pack(side="left", padx=10)
        ctk.CTkEntry(frame1, textvariable=self.ppt_path, width=380).pack(side="left")
        ctk.CTkButton(frame1, text="浏览选择", command=self.select_ppt, width=90).pack(side="left", padx=8)

        frame2 = ctk.CTkFrame(self)
        frame2.pack(padx=25, pady=8, fill="x")
        ctk.CTkLabel(frame2, text="输出视频名称：", font=("微软雅黑", 12)).pack(side="left", padx=10)
        ctk.CTkEntry(frame2, textvariable=self.output_name, width=380).pack(side="left")
        ctk.CTkLabel(frame2, text=".mp4", font=("微软雅黑", 12)).pack(side="left", padx=5)

        frame3 = ctk.CTkFrame(self)
        frame3.pack(padx=25, pady=8, fill="x")
        ctk.CTkLabel(frame3, text="字幕字体大小：", font=("微软雅黑", 12)).pack(side="left", padx=10)
        font_size_entry = ctk.CTkEntry(frame3, textvariable=self.subtitle_font_size, width=100)
        font_size_entry.pack(side="left")
        ctk.CTkLabel(frame3, text="号 (12-40)", font=("微软雅黑", 12)).pack(side="left", padx=5)

        self.progress_bar = ctk.CTkProgressBar(self)
        self.progress_bar.pack(padx=25, pady=5, fill="x")
        self.progress_bar.set(0)

        self.log_text = ctk.CTkTextbox(self, height=130, font=("微软雅黑", 10))
        self.log_text.pack(padx=25, pady=10, fill="both", expand=True)

        btn_frame = ctk.CTkFrame(self)
        btn_frame.pack(padx=25, pady=10, fill="x")
        
        self.convert_btn = ctk.CTkButton(
            btn_frame, text="开始转换生成视频", command=self.task_start,
            font=("微软雅黑", 14, "bold"), height=48
        )
        self.convert_btn.pack(side="left", fill="x", expand=True, padx=(0,5))
        
        self.stop_btn = ctk.CTkButton(
            btn_frame, text="强制停止", command=self.stop_conversion,
            font=("微软雅黑", 12, "bold"), height=48, state="disabled",
            fg_color="#ff4444", hover_color="#ff6666"
        )
        self.stop_btn.pack(side="right", padx=(5,0))

    def log(self, msg):
        try:
            self.log_text.insert("end", msg + "\n")
            self.log_text.see("end")
            self.update_idletasks()
        except:
            pass

    def select_ppt(self):
        path = filedialog.askopenfilename(
            title="选择PPT演示文稿",
            filetypes=[("PPT文件", "*.ppt;*.pptx"), ("所有文件", "*.*")]
        )
        if path:
            self.ppt_path.set(path)
            ppt_name = os.path.splitext(os.path.basename(path))[0]
            self.output_name.set(ppt_name)
            self.log(f"✅ 已选中文件：{os.path.basename(path)}")
            self.log(f"📝 输出视频名称已自动设置为：{ppt_name}")

    def stop_conversion(self):
        self.stop_requested = True
        self.log("⚠️ 强制停止已触发！")
        self.stop_btn.configure(state="disabled", text="正在停止...")

    def check_stop(self):
        if self.stop_requested:
            raise Exception("用户强制停止")

    def text_to_speech_offline(self, text, output_path):
        try:
            import win32com.client as win32
            
            speak = win32.Dispatch("SAPI.SpVoice")
            speak.Rate = 0
            speak.Volume = 100
            
            audio = win32.Dispatch("SAPI.SpFileStream")
            audio.Open(output_path, 3)
            speak.AudioOutputStream = audio
            
            speak.Speak(text)
            
            audio.Close()
            
            self.log("   🎤 使用系统SAPI离线语音合成")
            return True
        except Exception as e:
            self.log(f"   ⚠️  离线语音合成失败：{e}")
            return False

    async def text_to_speech_async(self, text, output_path, voice="zh-CN-XiaoxiaoNeural"):
        max_retries = 8
        for attempt in range(max_retries):
            try:
                self.log(f"   🌐 在线语音合成中... (尝试 {attempt+1}/{max_retries})")
                communicate = edge_tts.Communicate(text, voice)
                await communicate.save(output_path)
                self.log("   🎤 使用微软在线语音合成成功")
                return True
            except Exception as e:
                if attempt < max_retries - 1:
                    wait_time = 4 - (attempt // 2) if attempt < 6 else 1
                    self.log(f"   ⏳ 网络连接失败，{wait_time}秒后重试...")
                    await asyncio.sleep(wait_time)  # 重试间隔递减
                else:
                    self.log("   🌐 多次尝试失败，切换到离线语音...")
                    
                    if self.text_to_speech_offline(text, output_path):
                        return True
                    else:
                        self.log("   🔇 离线语音也失败，使用静默音频")
                        self.create_silent_audio(output_path, duration=3)
                        return True

    def text_to_speech(self, text, output_path):
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            result = loop.run_until_complete(self.text_to_speech_async(text, output_path))
            loop.close()
            return result
        except Exception as e:
            self.log(f"   ⚠️  语音合成异常：{e}")
            self.create_silent_audio(output_path, duration=3)
            return True

    def create_silent_audio(self, output_path, duration=3):
        try:
            import wave
            import struct
            sample_rate = 44100
            num_samples = duration * sample_rate
            
            with wave.open(output_path, 'w') as wav_file:
                wav_file.setnchannels(1)
                wav_file.setsampwidth(2)
                wav_file.setframerate(sample_rate)
                
                for _ in range(num_samples):
                    wav_file.writeframes(struct.pack('h', 0))
        except:
            pass

    def ppt_export_image(self, ppt_file):
        self.log("\n[1/3] 调用Office导出PPT高清图片...")
        ppt_app = None
        pres = None
        created_new = False
        
        try:
            self.check_stop()
            
            for img_file in (self.temp / "images").glob("*.png"):
                try: img_file.unlink()
                except: pass
            
            self.check_stop()
            
            try:
                ppt_app = win32.GetObject(Class="PowerPoint.Application")
                self.log("   检测到已运行的PowerPoint...")
                created_new = False
            except:
                ppt_app = win32.Dispatch("PowerPoint.Application")
                self.log("   启动新的PowerPoint实例...")
                created_new = True
            
            self.check_stop()
            
            try:
                ppt_app.DisplayAlerts = False
            except:
                pass
            
            self.check_stop()
            
            self.log("   正在打开PPT文件...")
            pres = ppt_app.Presentations.Open(os.path.abspath(ppt_file))
            
            self.check_stop()
            
            img_list = []
            total_slides = len(pres.Slides)
            self.log(f"   共{total_slides}页幻灯片")
            
            for idx, slide in enumerate(pres.Slides):
                self.check_stop()
                
                img_save = str(self.temp / "images" / f"slide_{idx:03d}.png")
                slide.Export(img_save, "PNG", 1920, 1080)
                img_list.append(img_save)
                self.log(f"   第{idx+1}/{total_slides}页导出完成")
                
                progress = (idx + 1) / (total_slides * 3)
                self.progress_bar.set(progress)
                self.update_idletasks()
            
            self.check_stop()
            
            self.log("   正在关闭PPT...")
            pres.Close()
            pres = None
            
            self.check_stop()
            
            if created_new:
                self.log("   正在关闭PowerPoint...")
                ppt_app.Quit()
                ppt_app = None
            
            self.log(f"   ✅ 图片导出完成，共{len(img_list)}张")
            return img_list
        except Exception as e:
            if str(e) != "用户强制停止":
                self.log(f"❌ PPT导出失败：{e}")
            if pres:
                try: pres.Close()
                except: pass
            if ppt_app and created_new:
                try: ppt_app.Quit()
                except: pass
            raise

    def generate_audios_from_notes(self, ppt_file):
        self.log("\n[2/3] 读取备注并生成语音...")
        pres = Presentation(ppt_file)
        audio_list = []
        note_texts = []
        
        for audio_file in (self.temp / "audios").glob("*.mp3"):
            try: audio_file.unlink()
            except: pass
        
        total_slides = len(pres.slides)
        self.log(f"   共{total_slides}页幻灯片待处理")
        
        for idx, slide in enumerate(pres.slides):
            self.check_stop()
            
            note_text = ""
            if slide.has_notes_slide:
                try:
                    note_text = slide.notes_slide.notes_text_frame.text.strip()
                except:
                    note_text = ""
            
            if not note_text:
                note_text = "本幻灯片无备注内容。"
            
            audio_save = str(self.temp / "audios" / f"voice_{idx:03d}.mp3")
            
            self.log(f"   第{idx+1}/{total_slides}页正在生成语音...")
            self.text_to_speech(note_text, audio_save)
            
            audio_list.append(audio_save)
            note_texts.append(note_text)
            self.log(f"   第{idx+1}/{total_slides}页语音合成完成")
            
            progress = (total_slides + idx + 1) / (total_slides * 3)
            self.progress_bar.set(progress)
            self.update_idletasks()
            
            time.sleep(0.1)
        
        self.log(f"   ✅ 语音生成完成，共{len(audio_list)}个")
        return audio_list, note_texts

    def format_srt_time(self, seconds):
        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        secs = int(seconds % 60)
        millis = int((seconds % 1) * 1000)
        return f"{hours:02d}:{minutes:02d}:{secs:02d},{millis:03d}"
    
    def draw_subtitle(self, img_path, text, output_path, font_size=20):
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            self.log(f"   🎨 正在绘制字幕到图片（{font_size}号）：{os.path.basename(output_path)}")
            
            img = Image.open(img_path)
            
            # 如果图片是RGBA模式，转成RGB
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            
            draw = ImageDraw.Draw(img, 'RGBA')
            
            width, height = img.size
            
            # 优化的自动换行 - 更长的行
            wrapped_text = []
            current_line = ""
            for char in text:
                current_line += char
                if len(current_line) >= 40 and char in "，。！？、；：":
                    wrapped_text.append(current_line)
                    current_line = ""
                    if len(wrapped_text) >= 2:
                        break
            if current_line and len(wrapped_text) < 2:
                wrapped_text.append(current_line)
            
            if not wrapped_text:
                wrapped_text = [text]
            
            # 字体设置
            font = None
            try:
                font = ImageFont.truetype("msyh.ttc", font_size)  # 微软雅黑
            except:
                try:
                    font = ImageFont.truetype("simhei.ttf", font_size)  # 黑体
                except:
                    try:
                        font = ImageFont.truetype("arial.ttf", font_size)
                    except:
                        font = ImageFont.load_default()
            
            # 计算文字尺寸
            line_height = 24
            total_height = len(wrapped_text) * line_height
            
            # 背景位置 - 非常靠下
            bg_y = height - total_height - 20
            
            # 绘制半透明背景
            bg_padding = 15
            bg_width = width - bg_padding * 2
            bg_height = total_height + 10
            
            draw.rectangle(
                [bg_padding, bg_y - 5, bg_padding + bg_width, bg_y + bg_height],
                fill=(0, 0, 0, 180)
            )
            
            # 绘制文字 - 用白色
            for i, line in enumerate(wrapped_text):
                try:
                    # 计算文字宽度居中
                    bbox = draw.textbbox((0, 0), line, font=font)
                    text_w = bbox[2] - bbox[0]
                    text_x = (width - text_w) // 2
                    text_y = bg_y + i * line_height
                    
                    draw.text((text_x, text_y), line, fill=(255, 255, 255, 255), font=font)
                except:
                    # 如果绘制失败，用简单方式
                    draw.text((width//4, bg_y + i*line_height), line, fill=(255, 255, 255, 255))
            
            img.save(output_path)
            self.log(f"   ✅ 字幕绘制成功：{os.path.basename(output_path)}")
            return True
        except Exception as e:
            self.log(f"   ❌ 绘制字幕失败：{e}")
            # 出错时直接复制原图
            try:
                import shutil
                shutil.copy(img_path, output_path)
                self.log(f"   ⚠️ 使用原图，没有字幕")
            except:
                pass
            return False

    def generate_subtitles(self, audio_paths, note_texts, srt_path):
        self.log("   📝 正在生成字幕文件...")
        current_time = 0.0
        with open(srt_path, 'w', encoding='utf-8') as f:
            for idx, (audio_path, note_text) in enumerate(zip(audio_paths, note_texts)):
                try:
                    aud_clip = AudioFileClip(audio_path)
                    duration = aud_clip.duration
                    aud_clip.close()
                    
                    start_time = current_time
                    end_time = current_time + duration
                    
                    f.write(f"{idx + 1}\n")
                    f.write(f"{self.format_srt_time(start_time)} --> {self.format_srt_time(end_time)}\n")
                    f.write(f"{note_text}\n\n")
                    
                    current_time = end_time
                except Exception as e:
                    continue
        
        self.log(f"   ✅ 字幕文件已生成：{srt_path}")

    def merge_video_audio(self, img_paths, audio_paths, note_texts, out_file, font_size=20):
        self.log("\n[3/3] 合成带配音视频...")
        clips = []
        total = len(img_paths)
        
        try:
            for idx, (img, aud, note) in enumerate(zip(img_paths, audio_paths, note_texts)):
                self.check_stop()
                
                if not os.path.exists(img) or not os.path.exists(aud):
                    self.log(f"   ⚠️  跳过缺少文件的第{idx+1}页")
                    continue
                
                try:
                    aud_clip = AudioFileClip(aud)
                    duration = aud_clip.duration
                    
                    # 用PIL绘制字幕
                    final_img_path = os.path.join(self.temp, f'slide_{idx:03d}_with_subtitle.png')
                    self.draw_subtitle(img, note, final_img_path, font_size)
                    
                    img_clip = ImageClip(final_img_path)
                    img_clip.duration = duration
                    img_clip.audio = aud_clip
                    
                    clips.append(img_clip)
                    
                    progress = (total * 2 + idx + 1) / (total * 3)
                    self.progress_bar.set(progress)
                    self.update_idletasks()
                    
                    if (idx + 1) % 5 == 0 or idx + 1 == total:
                        self.log(f"   处理进度：{idx+1}/{total}")
                except Exception as e:
                    self.log(f"   ⚠️  第{idx+1}页处理出错：{e}")
                    continue
            
            self.check_stop()
            
            if not clips:
                raise Exception("没有可用的幻灯片内容")
            
            self.log("   正在合成视频...")
            final_video = concatenate_videoclips(clips)
            
            self.check_stop()
            
            out_dir = os.path.dirname(os.path.abspath(out_file))
            os.makedirs(out_dir, exist_ok=True)
            
            self.log("   正在写入视频文件...")
            final_video.write_videofile(
                out_file,
                fps=24,
                codec="libx264",
                audio_codec="aac",
                preset="ultrafast",
                logger=None,
                threads=4
            )
            
            self.progress_bar.set(1)
            self.log("   ✅ 视频合成完成")
            
        finally:
            for c in clips:
                try: c.close()
                except: pass
            try: final_video.close()
            except: pass

    def task_thread(self, ppt_path, out_name):
        try:
            self.log("🚀 开始转换...")
            
            # 获取字体大小，确保在合理范围内
            font_size = self.subtitle_font_size.get()
            font_size = max(12, min(40, font_size))  # 限制在12-40之间
            self.log(f"📝 使用字幕字体大小：{font_size}号")
            
            imgs = self.ppt_export_image(ppt_path)
            
            audios, note_texts = self.generate_audios_from_notes(ppt_path)
            
            srt_path = os.path.splitext(out_name)[0] + ".srt"
            self.generate_subtitles(audios, note_texts, srt_path)
            
            self.merge_video_audio(imgs, audios, note_texts, out_name, font_size)
            
            self.after(0, lambda: self.on_conversion_success(out_name, srt_path))
        except Exception as e:
            if str(e) == "用户强制停止":
                self.log("⏹️  转换已被用户停止")
                self.after(0, self.on_conversion_stopped)
            else:
                self.after(0, lambda: self.on_conversion_error(str(e)))
        finally:
            self.after(0, self.on_conversion_finish)

    def on_conversion_success(self, out_name, srt_path):
        self.log(f"\n🎉 全部完成！")
        self.log(f"   📽️  视频已保存：{out_name}")
        self.log(f"   📝  字幕已保存：{srt_path}")
        messagebox.showinfo("完成", f"转换成功！\n视频：{os.path.basename(out_name)}\n字幕：{os.path.basename(srt_path)}")
        try:
            os.startfile(os.path.dirname(os.path.abspath(out_name)))
        except:
            pass

    def on_conversion_error(self, error_msg):
        self.log(f"\n❌ 转换失败：{error_msg}")
        messagebox.showerror("失败", f"运行错误：{error_msg}")

    def on_conversion_stopped(self):
        messagebox.showinfo("已停止", "转换已停止！")

    def on_conversion_finish(self):
        self.is_converting = False
        self.stop_requested = False
        self.convert_btn.configure(state="normal", text="开始转换生成视频")
        self.stop_btn.configure(state="disabled", text="强制停止")

    def task_start(self):
        if self.is_converting:
            messagebox.showwarning("警告", "正在转换中，请稍候...")
            return
        
        ppt_path = self.ppt_path.get().strip()
        out_name = self.output_name.get().strip() + ".mp4"
        
        if not os.path.exists(ppt_path):
            messagebox.showerror("错误", "请先选择有效的PPT/PPTX文件！")
            return
        
        self.is_converting = True
        self.stop_requested = False
        self.convert_btn.configure(state="disabled", text="正在转换...")
        self.stop_btn.configure(state="normal")
        self.progress_bar.set(0)
        
        thread = threading.Thread(target=self.task_thread, args=(ppt_path, out_name))
        thread.daemon = True
        thread.start()


if __name__ == "__main__":
    app = PPTToVideoGUI()
    app.mainloop()
